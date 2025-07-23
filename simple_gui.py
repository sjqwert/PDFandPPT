#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
简化版PDF转PPT GUI界面
不依赖PyMuPDF，避免macOS版本兼容性问题
"""

import os
import sys
import time
import logging
import threading
import tkinter as tk
from tkinter import ttk, filedialog, messagebox
from tkinter.scrolledtext import ScrolledText
from typing import List, Dict, Optional, Tuple, Any
import json
from pathlib import Path
from datetime import datetime
import subprocess
import importlib.util
import platform
import tempfile

# 设置日志
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger("simple_gui")

# 修补macOS版本检查
def patch_macos_version():
    """修补macOS版本检查，绕过PyMuPDF的版本限制"""
    if sys.platform == 'darwin':
        logger.info("检测到macOS系统，应用版本补丁")
        # 保存原始mac_ver函数
        original_mac_ver = platform.mac_ver
        
        # 创建一个修改版本的mac_ver函数，返回更高的macOS版本
        def patched_mac_ver():
            return ('13.7.0', ('', '', ''), 'x86_64')
        
        # 替换platform.mac_ver函数
        platform.mac_ver = patched_mac_ver
        logger.info("已应用macOS版本补丁")

# 应用补丁
patch_macos_version()

# 检查并安装必要的依赖
def check_and_install_dependencies():
    """检查并安装必要的依赖"""
    dependencies = {
        "pptx": "python-pptx"
    }
    
    missing_deps = []
    
    # 检查每个依赖
    for module_name, package_name in dependencies.items():
        if importlib.util.find_spec(module_name) is None:
            missing_deps.append(package_name)
    
    # 如果有缺失的依赖，尝试安装
    if missing_deps:
        try:
            logger.info(f"正在安装缺失的依赖: {', '.join(missing_deps)}")
            subprocess.check_call([sys.executable, "-m", "pip", "install"] + missing_deps)
            logger.info("依赖安装完成")
            return True
        except subprocess.CalledProcessError as e:
            logger.error(f"安装依赖失败: {str(e)}")
            messagebox.showerror("错误", f"无法安装必要的依赖。请尝试手动安装:\n\npip install {' '.join(missing_deps)}")
            sys.exit(1)
    
    return False

# 在导入其他模块前先检查依赖
check_and_install_dependencies()

# 导入python-pptx
try:
    from pptx import Presentation
    from pptx.util import Inches
    PPTX_AVAILABLE = True
    logger.info("成功导入python-pptx")
except ImportError:
    PPTX_AVAILABLE = False
    logger.error("无法导入python-pptx模块，程序无法运行")
    sys.exit(1)

# 自定义PDF转PPT函数
def convert_pdf_to_pptx(input_path, output_path, resolution=300):
    """使用外部工具将PDF转换为图片，然后创建PPT"""
    try:
        # 创建临时目录
        temp_dir = tempfile.mkdtemp()
        logger.info(f"创建临时目录: {temp_dir}")
        
        try:
            # 使用pdftoppm将PDF转换为图片
            cmd = [
                "pdftoppm",
                "-png",
                "-r", str(resolution),
                input_path,
                os.path.join(temp_dir, "page")
            ]
            logger.info(f"执行命令: {' '.join(cmd)}")
            result = subprocess.run(cmd, capture_output=True, text=True)
            
            if result.returncode != 0:
                logger.error(f"PDF转图片失败: {result.stderr}")
                return False
            
            # 获取生成的图片文件
            image_files = sorted([
                os.path.join(temp_dir, f)
                for f in os.listdir(temp_dir)
                if f.startswith("page-") and f.endswith(".png")
            ])
            
            if not image_files:
                logger.error("未生成任何图片文件")
                return False
            
            logger.info(f"生成了 {len(image_files)} 个图片文件")
            
            # 创建PPT
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            # 获取空白布局
            blank_layout = prs.slide_layouts[6]  # 通常索引6是空白布局
            
            # 添加图片到PPT
            for img_file in image_files:
                # 创建新的幻灯片
                slide = prs.slides.add_slide(blank_layout)
                
                # 添加图像到幻灯片
                slide.shapes.add_picture(
                    img_file, 
                    0, 0, 
                    width=prs.slide_width,
                    height=prs.slide_height
                )
            
            # 保存PPT文件
            prs.save(output_path)
            logger.info(f"已保存PPT文件: {output_path}")
            
            return True
        finally:
            # 清理临时文件
            import shutil
            shutil.rmtree(temp_dir, ignore_errors=True)
            logger.info(f"已清理临时目录: {temp_dir}")
    
    except Exception as e:
        logger.error(f"转换过程中出错: {str(e)}")
        return False

class SimplePDF2PPTApp:
    """简化版PDF转PPT应用"""
    
    def __init__(self, root):
        """初始化应用"""
        self.root = root
        self.root.title("简化版PDF转PPT工具")
        self.root.geometry("800x600")
        self.root.minsize(700, 500)
        
        # 窗口居中
        self.center_window()
        
        # 设置变量
        self.input_files = []  # 输入文件列表
        self.output_dir = tk.StringVar(value=os.path.expanduser("~"))  # 输出目录
        self.resolution = tk.IntVar(value=300)  # 分辨率
        
        # 创建UI
        self.create_ui()
        
        # 转换状态
        self.is_converting = False
        self.current_task = None
    
    def center_window(self):
        """使窗口在屏幕上居中显示"""
        # 更新窗口信息，确保获取正确的尺寸
        self.root.update_idletasks()
        
        # 获取屏幕宽度和高度
        screen_width = self.root.winfo_screenwidth()
        screen_height = self.root.winfo_screenheight()
        
        # 获取窗口宽度和高度
        window_width = self.root.winfo_width()
        window_height = self.root.winfo_height()
        
        # 计算居中位置的坐标
        x = (screen_width - window_width) // 2
        y = (screen_height - window_height) // 2
        
        # 设置窗口位置
        self.root.geometry(f"{window_width}x{window_height}+{x}+{y}")
    
    def create_ui(self):
        """创建用户界面"""
        # 创建主框架
        main_frame = ttk.Frame(self.root, padding=10)
        main_frame.pack(fill=tk.BOTH, expand=True)
        
        # 创建左右分栏
        left_frame = ttk.Frame(main_frame)
        left_frame.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        
        right_frame = ttk.Frame(main_frame)
        right_frame.pack(side=tk.RIGHT, fill=tk.BOTH, expand=True, padx=(10, 0))
        
        # 左侧：文件列表和操作按钮
        self.create_file_list_section(left_frame)
        
        # 右侧：转换选项和日志
        self.create_options_section(right_frame)
        self.create_log_section(right_frame)
        
        # 底部：状态栏
        self.create_status_bar()
    
    def create_file_list_section(self, parent):
        """创建文件列表部分"""
        # 文件列表框架
        file_frame = ttk.LabelFrame(parent, text="PDF文件列表", padding=10)
        file_frame.pack(fill=tk.BOTH, expand=True)
        
        # 文件列表
        file_list_frame = ttk.Frame(file_frame)
        file_list_frame.pack(fill=tk.BOTH, expand=True)
        
        # 创建带滚动条的列表框
        scrollbar = ttk.Scrollbar(file_list_frame)
        scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        
        self.file_listbox = tk.Listbox(file_list_frame, selectmode=tk.EXTENDED)
        self.file_listbox.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        self.file_listbox.config(yscrollcommand=scrollbar.set)
        scrollbar.config(command=self.file_listbox.yview)
        
        # 文件操作按钮
        btn_frame = ttk.Frame(file_frame, padding=(0, 10, 0, 0))
        btn_frame.pack(fill=tk.X)
        
        ttk.Button(btn_frame, text="添加文件", command=self.add_files).pack(side=tk.LEFT, padx=5)
        ttk.Button(btn_frame, text="添加文件夹", command=self.add_folder).pack(side=tk.LEFT, padx=5)
        ttk.Button(btn_frame, text="移除选中", command=self.remove_selected).pack(side=tk.LEFT, padx=5)
        ttk.Button(btn_frame, text="清空列表", command=self.clear_list).pack(side=tk.LEFT, padx=5)
        
        # 输出目录选择
        out_frame = ttk.Frame(file_frame, padding=(0, 10, 0, 0))
        out_frame.pack(fill=tk.X)
        
        ttk.Label(out_frame, text="输出目录:").pack(side=tk.LEFT)
        ttk.Entry(out_frame, textvariable=self.output_dir).pack(side=tk.LEFT, fill=tk.X, expand=True, padx=5)
        ttk.Button(out_frame, text="浏览...", command=self.browse_output_dir).pack(side=tk.LEFT)
        
        # 转换按钮
        convert_frame = ttk.Frame(file_frame, padding=(0, 10, 0, 0))
        convert_frame.pack(fill=tk.X)
        
        self.convert_btn = ttk.Button(convert_frame, text="开始转换", command=self.start_conversion)
        self.convert_btn.pack(fill=tk.X, pady=5)
        
        # 进度条
        self.progress_var = tk.DoubleVar(value=0.0)
        self.progress = ttk.Progressbar(file_frame, variable=self.progress_var, maximum=100)
        self.progress.pack(fill=tk.X, pady=5)
    
    def create_options_section(self, parent):
        """创建转换选项部分"""
        options_frame = ttk.LabelFrame(parent, text="转换选项", padding=10)
        options_frame.pack(fill=tk.X)
        
        # 分辨率
        res_frame = ttk.Frame(options_frame)
        res_frame.pack(fill=tk.X, pady=5)
        
        ttk.Label(res_frame, text="分辨率(DPI):").pack(side=tk.LEFT)
        res_spinbox = ttk.Spinbox(res_frame, from_=72, to=600, increment=1, textvariable=self.resolution)
        res_spinbox.pack(side=tk.LEFT, padx=5, fill=tk.X, expand=True)
    
    def create_log_section(self, parent):
        """创建日志部分"""
        log_frame = ttk.LabelFrame(parent, text="转换日志", padding=10)
        log_frame.pack(fill=tk.BOTH, expand=True, pady=10)
        
        # 日志文本框
        self.log_text = ScrolledText(log_frame, wrap=tk.WORD, height=10)
        self.log_text.pack(fill=tk.BOTH, expand=True)
        self.log_text.config(state=tk.DISABLED)
        
        # 添加日志处理器
        self.log_handler = TextHandler(self.log_text)
        self.log_handler.setLevel(logging.INFO)
        logger.addHandler(self.log_handler)
        
        # 清除日志按钮
        ttk.Button(log_frame, text="清除日志", command=self.clear_log).pack(anchor=tk.E, pady=(5, 0))
    
    def create_status_bar(self):
        """创建状态栏"""
        self.status_bar = ttk.Label(self.root, text="就绪", relief=tk.SUNKEN, anchor=tk.W)
        self.status_bar.pack(side=tk.BOTTOM, fill=tk.X)
    
    def add_files(self):
        """添加PDF文件"""
        files = filedialog.askopenfilenames(
            title="选择PDF文件",
            filetypes=[("PDF文件", "*.pdf"), ("所有文件", "*.*")]
        )
        
        if files:
            for file in files:
                if file.lower().endswith('.pdf') and file not in self.input_files:
                    self.input_files.append(file)
                    self.file_listbox.insert(tk.END, os.path.basename(file))
            
            self.update_status()
    
    def add_folder(self):
        """添加文件夹中的所有PDF文件"""
        folder = filedialog.askdirectory(title="选择包含PDF文件的文件夹")
        
        if folder:
            pdf_files = []
            for root, _, files in os.walk(folder):
                for file in files:
                    if file.lower().endswith(".pdf"):
                        pdf_path = os.path.join(root, file)
                        if pdf_path not in self.input_files:
                            pdf_files.append(pdf_path)
            
            if pdf_files:
                for file in pdf_files:
                    self.input_files.append(file)
                    self.file_listbox.insert(tk.END, os.path.basename(file))
                
                messagebox.showinfo("添加成功", f"已添加 {len(pdf_files)} 个PDF文件")
                self.update_status()
            else:
                messagebox.showinfo("提示", "所选文件夹中没有找到PDF文件")
    
    def remove_selected(self):
        """移除选中的文件"""
        selected_indices = self.file_listbox.curselection()
        
        if not selected_indices:
            return
        
        # 从后往前删除，避免索引变化
        for i in sorted(selected_indices, reverse=True):
            del self.input_files[i]
            self.file_listbox.delete(i)
        
        self.update_status()
    
    def clear_list(self):
        """清空文件列表"""
        self.input_files = []
        self.file_listbox.delete(0, tk.END)
        self.update_status()
    
    def browse_output_dir(self):
        """浏览输出目录"""
        directory = filedialog.askdirectory(title="选择输出目录")
        
        if directory:
            self.output_dir.set(directory)
    
    def update_status(self):
        """更新状态栏信息"""
        status_text = f"文件数: {len(self.input_files)}"
        
        # 添加当前状态
        if self.is_converting:
            status_text = "正在转换... | " + status_text
        
        self.status_bar.config(text=status_text)
    
    def start_conversion(self):
        """开始转换过程"""
        if not self.input_files:
            messagebox.showwarning("警告", "请先添加PDF文件")
            return
        
        if self.is_converting:
            messagebox.showinfo("提示", "正在进行转换，请等待当前任务完成")
            return
        
        # 设置转换状态
        self.is_converting = True
        self.convert_btn.config(text="转换中...", state=tk.DISABLED)
        self.update_status()
        
        # 获取转换选项
        output_dir = self.output_dir.get()
        resolution = self.resolution.get()
        
        # 重置进度条
        self.progress_var.set(0)
        
        # 在新线程中执行转换
        self.current_task = threading.Thread(
            target=self.convert_files,
            args=(output_dir, resolution)
        )
        self.current_task.daemon = True
        self.current_task.start()
    
    def convert_files(self, output_dir, resolution):
        """在后台线程中转换文件"""
        try:
            total_files = len(self.input_files)
            successful = 0
            failed = 0
            
            logger.info(f"开始批量转换 {total_files} 个文件")
            logger.info(f"分辨率: {resolution} DPI")
            
            start_time = time.time()
            
            for i, input_file in enumerate(self.input_files):
                # 更新进度
                progress = (i / total_files) * 100
                self.progress_var.set(progress)
                self.root.update_idletasks()
                
                # 设置输出路径
                file_name = os.path.basename(input_file)
                base_name = os.path.splitext(file_name)[0]
                output_file = os.path.join(output_dir, f"{base_name}.pptx")
                
                logger.info(f"正在转换 ({i+1}/{total_files}): {file_name}")
                
                try:
                    # 执行转换
                    convert_pdf_to_pptx(
                        input_file,
                        output_file,
                        resolution=resolution
                    )
                    
                    # 检查输出文件是否生成
                    if os.path.exists(output_file):
                        logger.info(f"转换成功: {output_file}")
                        successful += 1
                    else:
                        logger.error(f"转换失败，未生成输出文件: {output_file}")
                        failed += 1
                except Exception as e:
                    logger.error(f"转换失败: {str(e)}")
                    failed += 1
            
            # 完成转换
            elapsed_time = time.time() - start_time
            logger.info(f"批量转换完成，耗时: {elapsed_time:.2f}秒")
            logger.info(f"成功: {successful}，失败: {failed}")
            
            # 设置进度为100%
            self.progress_var.set(100)
            
            # 显示完成消息
            messagebox.showinfo("转换完成", f"成功转换 {successful} 个文件，失败 {failed} 个文件")
        
        except Exception as e:
            logger.error(f"转换过程中出错: {str(e)}")
            messagebox.showerror("错误", f"转换过程中出错: {str(e)}")
        
        finally:
            # 重置转换状态
            self.is_converting = False
            self.convert_btn.config(text="开始转换", state=tk.NORMAL)
            self.update_status()
    
    def clear_log(self):
        """清除日志"""
        self.log_text.config(state=tk.NORMAL)
        self.log_text.delete(1.0, tk.END)
        self.log_text.config(state=tk.DISABLED)

class TextHandler(logging.Handler):
    """将日志输出到Tkinter文本框的处理器"""
    
    def __init__(self, text_widget):
        logging.Handler.__init__(self)
        self.text_widget = text_widget
    
    def emit(self, record):
        msg = self.format(record)
        
        def append():
            self.text_widget.config(state=tk.NORMAL)
            self.text_widget.insert(tk.END, msg + "\n")
            self.text_widget.see(tk.END)
            self.text_widget.config(state=tk.DISABLED)
        
        # 在主线程中更新UI
        self.text_widget.after(0, append)

def main():
    """主函数"""
    root = tk.Tk()
    app = SimplePDF2PPTApp(root)
    
    # 设置图标（如果有）
    icon_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "pdf2pptx.ico")
    if os.path.exists(icon_path):
        try:
            root.iconbitmap(icon_path)
        except tk.TclError:
            logger.warning("无法设置应用图标")
    
    root.mainloop()

if __name__ == "__main__":
    main() 