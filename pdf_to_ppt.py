#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import os
import sys
import subprocess
import tempfile
import logging
from pathlib import Path
from typing import Optional, Tuple, Dict, Any, List, Union, cast
from enum import Enum

# 设置日志
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger("pdf_to_ppt")

# 尝试导入必要的库
try:
    from pdf2pptx import convert_pdf2pptx
    PDF2PPTX_AVAILABLE = True
except ImportError:
    PDF2PPTX_AVAILABLE = False
    logger.warning("pdf2pptx库未安装，将使用替代方法进行转换")

# 尝试导入PyMuPDF，使用monkey patch绕过macOS版本检查
try:
    # 首先尝试导入sys模块
    import sys
    import platform
    
    # 保存原始platform.mac_ver函数
    original_mac_ver = platform.mac_ver
    
    # 创建一个修改版本的mac_ver函数，返回更高的macOS版本
    def patched_mac_ver():
        return ('13.7.0', ('', '', ''), 'x86_64')
    
    # 替换platform.mac_ver函数
    platform.mac_ver = patched_mac_ver
    
    # 现在尝试导入fitz
    import fitz
    
    # 导入成功后，恢复原始函数
    platform.mac_ver = original_mac_ver
    
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False
    logger.warning("PyMuPDF (fitz) 库未安装，某些功能将不可用")

try:
    from pptx import Presentation
    from pptx.util import Inches
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx库未安装，将无法使用替代转换方法")

class ConversionMethod(Enum):
    """转换方法枚举"""
    PDF2PPTX = "pdf2pptx"  # 使用pdf2pptx库
    PYMUPDF = "pymupdf"    # 使用PyMuPDF + python-pptx
    IMAGEMAGICK = "imagemagick"  # 使用ImageMagick + python-pptx

class PDFAspectRatio(Enum):
    """PDF纵横比枚举"""
    STANDARD_4_3 = "4:3"  # 标准4:3比例
    WIDESCREEN_16_9 = "16:9"  # 宽屏16:9比例
    CUSTOM = "custom"  # 自定义比例

class PDFtoPPTConverter:
    """增强版PDF转PPT转换器类"""
    
    @staticmethod
    def check_dependencies() -> Dict[str, bool]:
        """
        检查所有依赖是否满足
        
        返回:
            包含各依赖可用性的字典
        """
        return {
            "pdf2pptx": PDF2PPTX_AVAILABLE,
            "pymupdf": PYMUPDF_AVAILABLE,
            "pptx": PPTX_AVAILABLE,
            "imagemagick": PDFtoPPTConverter._check_imagemagick()
        }
    
    @staticmethod
    def _check_imagemagick() -> bool:
        """检查ImageMagick是否可用"""
        try:
            result = subprocess.run(["convert", "--version"], 
                                   capture_output=True, text=True)
            return result.returncode == 0 and "ImageMagick" in result.stdout
        except (FileNotFoundError, subprocess.SubprocessError):
            return False
    
    @staticmethod
    def detect_aspect_ratio(pdf_path: str) -> Tuple[PDFAspectRatio, float]:
        """
        检测PDF的纵横比
        
        参数:
            pdf_path: PDF文件路径
            
        返回:
            (纵横比类型, 具体比值)
        """
        if PYMUPDF_AVAILABLE:
            try:
                doc = fitz.open(pdf_path)
                page = doc[0]  # 获取第一页
                width, height = page.rect.width, page.rect.height
                ratio = height / width
                doc.close()
                
                # 判断纵横比类型
                if 0.74 <= ratio <= 0.76:  # 4:3 比例 (0.75)
                    return PDFAspectRatio.STANDARD_4_3, ratio
                elif 0.55 <= ratio <= 0.57:  # 16:9 比例 (0.5625)
                    return PDFAspectRatio.WIDESCREEN_16_9, ratio
                else:
                    return PDFAspectRatio.CUSTOM, ratio
            except Exception as e:
                logger.error(f"检测PDF纵横比时出错: {str(e)}")
                return PDFAspectRatio.STANDARD_4_3, 0.75
        else:
            # 如果PyMuPDF不可用，默认使用4:3比例
            logger.warning("PyMuPDF不可用，无法检测PDF纵横比，将使用默认4:3比例")
            return PDFAspectRatio.STANDARD_4_3, 0.75
    
    @staticmethod
    def convert_pdf_to_ppt(
        input_path: str, 
        output_path: Optional[str] = None,
        method: ConversionMethod = ConversionMethod.PDF2PPTX,
        resolution: int = 300,
        start_page: int = 0,
        page_count: Optional[int] = None,
        template_path: Optional[str] = None,
        detect_ratio: bool = True
    ) -> Tuple[bool, str]:
        """
        将PDF转换为PPT
        
        参数:
            input_path: 输入PDF文件路径
            output_path: 输出PPT文件路径，如果为None则自动生成
            method: 转换方法
            resolution: 图像分辨率
            start_page: 起始页码（从0开始）
            page_count: 要转换的页数，None表示全部
            template_path: PPT模板路径
            detect_ratio: 是否自动检测并适配PDF纵横比
            
        返回:
            (成功标志, 输出文件路径或错误信息)
        """
        logger.info(f"开始转换PDF到PPT: {input_path}")
        logger.info(f"使用转换方法: {method.value}")
        
        # 检查输入文件是否存在
        if not os.path.exists(input_path):
            return False, f"输入文件不存在: {input_path}"
        
        # 如果没有指定输出路径，则自动生成
        if output_path is None:
            input_file = Path(input_path)
            output_path = str(input_file.parent / f"{input_file.stem}.pptx")
        
        # 根据选择的方法进行转换
        if method == ConversionMethod.PDF2PPTX:
            return PDFtoPPTConverter._convert_with_pdf2pptx(
                input_path, output_path, resolution, start_page, page_count, detect_ratio
            )
        elif method == ConversionMethod.PYMUPDF:
            return PDFtoPPTConverter._convert_with_pymupdf(
                input_path, output_path, resolution, start_page, page_count, 
                template_path, detect_ratio
            )
        elif method == ConversionMethod.IMAGEMAGICK:
            return PDFtoPPTConverter._convert_with_imagemagick(
                input_path, output_path, resolution, start_page, page_count,
                template_path, detect_ratio
            )
        else:
            return False, f"不支持的转换方法: {method}"
    
    @staticmethod
    def _convert_with_pdf2pptx(
        input_path: str, 
        output_path: str,
        resolution: int = 300,
        start_page: int = 0,
        page_count: Optional[int] = None,
        detect_ratio: bool = True
    ) -> Tuple[bool, str]:
        """使用pdf2pptx库转换"""
        if not PDF2PPTX_AVAILABLE:
            return False, "缺少必要的库: pdf2pptx。请运行 'pip install pdf2pptx python-pptx' 安装。"
        
        try:
            logger.info(f"使用pdf2pptx库转换: {input_path} -> {output_path}")
            
            # 执行转换，添加必要的参数
            convert_pdf2pptx(
                input_path,
                output_path,
                resolution=resolution,
                start_page=start_page,
                page_count=page_count
            )
            
            # 检查输出文件是否生成
            if os.path.exists(output_path):
                logger.info(f"转换成功: {output_path}")
                return True, output_path
            else:
                logger.error("转换失败，未生成输出文件")
                return False, "转换失败，未生成输出文件"
            
        except Exception as e:
            logger.error(f"pdf2pptx转换过程中出错: {str(e)}")
            return False, f"转换过程中出错: {str(e)}"
    
    @staticmethod
    def _convert_with_pymupdf(
        input_path: str, 
        output_path: str,
        resolution: int = 300,
        start_page: int = 0,
        page_count: Optional[int] = None,
        template_path: Optional[str] = None,
        detect_ratio: bool = True
    ) -> Tuple[bool, str]:
        """使用PyMuPDF + python-pptx转换"""
        if not PYMUPDF_AVAILABLE:
            return False, "缺少必要的库: PyMuPDF。请运行 'pip install pymupdf' 安装。"
        
        if not PPTX_AVAILABLE:
            return False, "缺少必要的库: python-pptx。请运行 'pip install python-pptx' 安装。"
        
        try:
            logger.info(f"使用PyMuPDF转换: {input_path} -> {output_path}")
            
            # 打开PDF文件
            pdf = fitz.open(input_path)
            
            # 确定要处理的页面范围
            total_pages = pdf.page_count
            if page_count is None:
                end_page = total_pages
            else:
                end_page = min(start_page + page_count, total_pages)
            
            # 创建一个新的PPT演示文稿或使用模板
            if template_path and os.path.exists(template_path):
                prs = Presentation(template_path)
            else:
                prs = Presentation()
            
            # 如果需要检测并适配纵横比
            if detect_ratio:
                aspect_ratio, ratio_value = PDFtoPPTConverter.detect_aspect_ratio(input_path)
                if aspect_ratio == PDFAspectRatio.WIDESCREEN_16_9:
                    prs.slide_width = Inches(10)
                    prs.slide_height = Inches(5.625)
                else:  # 默认使用4:3比例
                    prs.slide_width = Inches(10)
                    prs.slide_height = Inches(7.5)
            
            # 获取空白布局
            blank_layout = prs.slide_layouts[6]  # 通常索引6是空白布局
            
            # 创建临时目录存储图片
            with tempfile.TemporaryDirectory() as temp_dir:
                # 遍历PDF页面
                for pg in range(start_page, end_page):
                    logger.info(f"处理第 {pg+1}/{total_pages} 页")
                    
                    # 获取页面
                    page = pdf[pg]
                    
                    # 设置渲染参数
                    zoom = resolution / 72  # 将DPI转换为缩放因子
                    matrix = fitz.Matrix(zoom, zoom)
                    
                    # 渲染页面为图像
                    pix = page.get_pixmap(matrix=matrix, alpha=False)
                    img_path = os.path.join(temp_dir, f"page_{pg+1:04d}.png")
                    pix.save(img_path)
                    
                    # 创建新的幻灯片
                    slide = prs.slides.add_slide(blank_layout)
                    
                    # 添加图像到幻灯片
                    slide.shapes.add_picture(
                        img_path, 
                        0, 0, 
                        width=prs.slide_width,
                        height=prs.slide_height
                    )
                
                # 保存PPT文件
                prs.save(output_path)
            
            # 关闭PDF文件
            pdf.close()
            
            # 检查输出文件是否生成
            if os.path.exists(output_path):
                logger.info(f"转换成功: {output_path}")
                return True, output_path
            else:
                logger.error("转换失败，未生成输出文件")
                return False, "转换失败，未生成输出文件"
            
        except Exception as e:
            logger.error(f"PyMuPDF转换过程中出错: {str(e)}")
            return False, f"转换过程中出错: {str(e)}"
    
    @staticmethod
    def _convert_with_imagemagick(
        input_path: str, 
        output_path: str,
        resolution: int = 300,
        start_page: int = 0,
        page_count: Optional[int] = None,
        template_path: Optional[str] = None,
        detect_ratio: bool = True
    ) -> Tuple[bool, str]:
        """使用ImageMagick + python-pptx转换"""
        if not PDFtoPPTConverter._check_imagemagick():
            return False, "ImageMagick未安装或不可用。请安装ImageMagick并确保'convert'命令可用。"
        
        if not PPTX_AVAILABLE:
            return False, "缺少必要的库: python-pptx。请运行 'pip install python-pptx' 安装。"
        
        try:
            logger.info(f"使用ImageMagick转换: {input_path} -> {output_path}")
            
            # 创建临时目录存储图片
            with tempfile.TemporaryDirectory() as temp_dir:
                # 构建ImageMagick命令
                img_cmd = [
                    "convert", 
                    "-density", str(resolution), 
                    input_path
                ]
                
                # 如果指定了页面范围
                if page_count is not None:
                    img_cmd.extend([
                        f"{input_path}[{start_page}-{start_page + page_count - 1}]"
                    ])
                
                # 输出路径
                img_cmd.append(f"{temp_dir}/page_%04d.png")
                
                # 执行转换命令
                logger.info(f"执行ImageMagick命令: {' '.join(img_cmd)}")
                result = subprocess.run(img_cmd, capture_output=True, text=True)
                
                if result.returncode != 0:
                    logger.error(f"ImageMagick转换失败: {result.stderr}")
                    return False, f"PDF转图片失败: {result.stderr}"
                
                # 创建一个新的PPT演示文稿或使用模板
                if template_path and os.path.exists(template_path):
                    prs = Presentation(template_path)
                else:
                    prs = Presentation()
                
                # 如果需要检测并适配纵横比
                if detect_ratio:
                    aspect_ratio, ratio_value = PDFtoPPTConverter.detect_aspect_ratio(input_path)
                    if aspect_ratio == PDFAspectRatio.WIDESCREEN_16_9:
                        prs.slide_width = Inches(10)
                        prs.slide_height = Inches(5.625)
                    else:  # 默认使用4:3比例
                        prs.slide_width = Inches(10)
                        prs.slide_height = Inches(7.5)
                
                # 获取空白布局
                blank_layout = prs.slide_layouts[6]  # 通常索引6是空白布局
                
                # 获取生成的图像文件
                image_files = sorted([
                    f for f in os.listdir(temp_dir) 
                    if f.startswith("page_") and f.endswith(".png")
                ])
                
                # 添加图像到PPT
                for img_file in image_files:
                    img_path = os.path.join(temp_dir, img_file)
                    
                    # 创建新的幻灯片
                    slide = prs.slides.add_slide(blank_layout)
                    
                    # 添加图像到幻灯片
                    slide.shapes.add_picture(
                        img_path, 
                        0, 0, 
                        width=prs.slide_width,
                        height=prs.slide_height
                    )
                
                # 保存PPT文件
                prs.save(output_path)
            
            # 检查输出文件是否生成
            if os.path.exists(output_path):
                logger.info(f"转换成功: {output_path}")
                return True, output_path
            else:
                logger.error("转换失败，未生成输出文件")
                return False, "转换失败，未生成输出文件"
            
        except Exception as e:
            logger.error(f"ImageMagick转换过程中出错: {str(e)}")
            return False, f"转换过程中出错: {str(e)}"
    
    @staticmethod
    def get_best_available_method() -> ConversionMethod:
        """获取最佳可用的转换方法"""
        deps = PDFtoPPTConverter.check_dependencies()
        
        if deps["pdf2pptx"]:
            return ConversionMethod.PDF2PPTX
        elif deps["pymupdf"] and deps["pptx"]:
            return ConversionMethod.PYMUPDF
        elif deps["imagemagick"] and deps["pptx"]:
            return ConversionMethod.IMAGEMAGICK
        else:
            logger.error("没有可用的转换方法")
            raise RuntimeError("没有可用的转换方法，请安装必要的依赖")

# 测试代码
if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("用法: python pdf_to_ppt.py <pdf文件路径> [ppt输出路径]")
        sys.exit(1)
    
    input_pdf = sys.argv[1]
    output_ppt = sys.argv[2] if len(sys.argv) > 2 else None
    
    # 获取最佳可用的转换方法
    try:
        method = PDFtoPPTConverter.get_best_available_method()
        print(f"使用转换方法: {method.value}")
        
        # 执行转换
        success, result = PDFtoPPTConverter.convert_pdf_to_ppt(
            input_pdf, 
            output_ppt,
            method=method,
            resolution=300,
            detect_ratio=True
        )
        
        if success:
            print(f"转换成功! 输出文件: {result}")
        else:
            print(f"转换失败: {result}")
    except Exception as e:
        print(f"转换过程中出错: {str(e)}")
        sys.exit(1) 